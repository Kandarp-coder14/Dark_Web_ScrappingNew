from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Dark Web Scraping for Cybersecurity and Threat Intelligence"
slide.placeholders[1].text = "Project Presentation\nDefensive and Legal Monitoring Approach"

# Slide 2: What is Dark Web Scraping?
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What is Dark Web Scraping?"
body = slide.shapes.placeholders[1].text_frame
body.text = "Collection of data from Tor-hosted (.onion) websites"
for t in [
    "Requires Tor browser/proxy instead of normal internet routing",
    "Useful for discovering early indicators of cyber threats",
    "Not indexed by traditional search engines"
]:
    p = body.add_paragraph()
    p.text = t

# Slide 3: Why Organizations Use It
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Why Organizations Use It"
body = slide.shapes.placeholders[1].text_frame
body.text = "Primary defensive use cases"
for t in [
    "Detect leaked credentials and exposed corporate data",
    "Track ransomware groups and threat actor chatter",
    "Identify attack trends and emerging TTPs early",
    "Support SOC, threat hunting, and incident response"
]:
    p = body.add_paragraph()
    p.text = t

# Slide 4: Project Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Project Architecture (Implemented)"
body = slide.shapes.placeholders[1].text_frame
body.text = "Pipeline modules in this project"
for t in [
    "Config Loader (YAML): target seeds, domain allowlist, limits",
    "Tor-aware HTTP Client: optional SOCKS proxy routing",
    "Crawler: controlled fetching with depth/page/time boundaries",
    "Parser: HTML extraction + indicator pattern matching",
    "Storage: SQLite tables for pages and indicators",
    "CLI Runner: executes end-to-end and prints run statistics"
]:
    p = body.add_paragraph()
    p.text = t

# Slide 5: Workflow
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "End-to-End Workflow"
body = slide.shapes.placeholders[1].text_frame
body.text = "1) Start from approved seed URLs"
for t in [
    "2) Fetch via Tor (if enabled) with request timeout",
    "3) Parse HTML, collect links, and extract indicators",
    "4) Enforce allowlist and crawl-depth restrictions",
    "5) Save output in SQLite for analyst review",
    "6) Generate metrics (visited, success, failed, indicators)"
]:
    p = body.add_paragraph()
    p.text = t

# Slide 6: Risks and Challenges
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Risks and Challenges"
body = slide.shapes.placeholders[1].text_frame
body.text = "Operational and legal concerns"
for t in [
    "Legal/ethical boundaries must be clearly defined",
    "Malware exposure risk from hostile sites",
    "Fake/unreliable data and disinformation",
    "Slow and unstable connectivity over Tor",
    "Need for secure, isolated execution environments"
]:
    p = body.add_paragraph()
    p.text = t

# Slide 7: Controls and Best Practices
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Controls and Best Practices"
body = slide.shapes.placeholders[1].text_frame
body.text = "Recommended safeguards"
for t in [
    "Use only authorized targets and documented legal approval",
    "Run in sandboxed VM/container with least privilege",
    "Maintain strict domain allowlist and crawl limits",
    "Store data securely and protect sensitive findings",
    "Review indicators before acting (human analyst validation)"
]:
    p = body.add_paragraph()
    p.text = t

# Slide 8: Project Demo Outcome
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Project Demo Outcome"
body = slide.shapes.placeholders[1].text_frame
body.text = "Implementation status in this repository"
for t in [
    "Dependencies installed successfully",
    "CLI execution completed with summary metrics",
    "SQLite database created at data/intel.db",
    "Ready for approved target configuration and reporting extensions"
]:
    p = body.add_paragraph()
    p.text = t

# Slide 9: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
body = slide.shapes.placeholders[1].text_frame
body.text = "Dark web scraping can improve cyber threat visibility when used responsibly"
for t in [
    "Strong value for threat intelligence and early warning",
    "Must be governed by legal, ethical, and operational controls",
    "This project provides a safe technical foundation for defensive research"
]:
    p = body.add_paragraph()
    p.text = t

out = "presentation/Dark_Web_Scraping_Threat_Intelligence_Project.pptx"
prs.save(out)
print(out)
